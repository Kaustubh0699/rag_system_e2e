"""Hierarchical recursive chunking with proper metadata for PDF and PPT."""

import logging
import re
from pathlib import Path
from typing import List, Tuple

import pypdf
import tiktoken
from pptx import Presentation

from .errors import DocumentParseError, UnsupportedFileFormatError
from .types import Chunk, Document

logger = logging.getLogger(__name__)

# Use cl100k_base encoding (GPT-4 style) for token counting
ENCODING = tiktoken.get_encoding("cl100k_base")


def count_tokens(text: str) -> int:
    """Count tokens in text."""
    return len(ENCODING.encode(text))


def split_sentences(text: str) -> List[str]:
    """Split text into sentences using regex."""
    # Split on sentence endings: . ! ? followed by space or end of string
    sentences = re.split(r'([.!?]+(?:\s+|$))', text)
    # Recombine sentences with their punctuation
    result = []
    for i in range(0, len(sentences) - 1, 2):
        if i + 1 < len(sentences):
            result.append(sentences[i] + sentences[i + 1])
        else:
            result.append(sentences[i])
    if len(sentences) % 2 == 1:
        result.append(sentences[-1])
    return [s.strip() for s in result if s.strip()]


def split_by_tokens(text: str, max_tokens: int) -> List[str]:
    """Soft split by tokens (fallback when sentence is too long)."""
    if count_tokens(text) <= max_tokens:
        return [text]
    
    words = text.split()
    chunks = []
    current_chunk = []
    current_tokens = 0
    
    for word in words:
        word_with_space = word + " "
        word_tokens = count_tokens(word_with_space)
        
        if current_tokens + word_tokens > max_tokens and current_chunk:
            chunks.append(" ".join(current_chunk))
            current_chunk = [word]
            current_tokens = count_tokens(word)
        else:
            current_chunk.append(word)
            current_tokens = count_tokens(" ".join(current_chunk))
    
    if current_chunk:
        chunks.append(" ".join(current_chunk))
    
    return chunks


def recursive_split(text: str, max_tokens: int, level: int = 0) -> List[str]:
    """Recursively split text using hierarchy: paragraphs -> lines -> sentences -> tokens."""
    if not text.strip():
        return []
    
    # Check if already fits
    if count_tokens(text) <= max_tokens:
        return [text]
    
    # Level 0: Split by paragraph breaks (\n\n)
    if level == 0:
        paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
        if len(paragraphs) > 1:
            result = []
            for para in paragraphs:
                result.extend(recursive_split(para, max_tokens, level + 1))
            return result
        else:
            return recursive_split(text, max_tokens, level + 1)
    
    # Level 1: Split by line breaks (\n)
    elif level == 1:
        lines = [l.strip() for l in text.split("\n") if l.strip()]
        if len(lines) > 1:
            result = []
            for line in lines:
                result.extend(recursive_split(line, max_tokens, level + 1))
            return result
        else:
            return recursive_split(text, max_tokens, level + 1)
    
    # Level 2: Split by sentences
    elif level == 2:
        sentences = split_sentences(text)
        if len(sentences) > 1:
            result = []
            for sentence in sentences:
                result.extend(recursive_split(sentence, max_tokens, level + 1))
            return result
        else:
            return recursive_split(text, max_tokens, level + 1)
    
    # Level 3: Fallback to token-based splitting
    else:
        return split_by_tokens(text, max_tokens)


def apply_sentence_overlap(chunks: List[str], overlap_sentences: int = 2) -> List[str]:
    """Apply sentence overlap between chunks.
    
    Takes the last N sentences from each chunk and prepends them to the next chunk.
    """
    if not chunks or overlap_sentences <= 0:
        return chunks
    
    result = []
    
    for i, chunk in enumerate(chunks):
        if i == 0:
            # First chunk - no overlap to add
            result.append(chunk)
        else:
            # Get last N sentences from previous chunk
            prev_chunk = chunks[i - 1]
            prev_sentences = split_sentences(prev_chunk)
            
            if len(prev_sentences) >= overlap_sentences:
                overlap_text = " ".join(prev_sentences[-overlap_sentences:])
                # Prepend overlap to current chunk
                overlapped_chunk = overlap_text + " " + chunk
                result.append(overlapped_chunk)
            else:
                # Not enough sentences for full overlap, use what we have
                if prev_sentences:
                    overlap_text = " ".join(prev_sentences)
                    overlapped_chunk = overlap_text + " " + chunk
                    result.append(overlapped_chunk)
                else:
                    result.append(chunk)
    
    return result


def parse_pdf(file_path: Path) -> List[Tuple[str, int]]:
    """Extract text from PDF file with page numbers.
    
    Returns:
        List of (text, page_num) tuples
    """
    try:
        pages = []
        with open(file_path, "rb") as f:
            reader = pypdf.PdfReader(f)
            for page_num, page in enumerate(reader.pages, start=1):
                text = page.extract_text()
                if text.strip():
                    pages.append((text, page_num))
        return pages
    except Exception as e:
        raise DocumentParseError(f"Failed to parse PDF: {e}") from e


def parse_ppt(file_path: Path) -> List[Tuple[str, dict]]:
    """Extract text from PPT/PPTX file with slide metadata.
    
    Note: Only .pptx format is supported (Office Open XML).
    Old .ppt format (binary) is not supported.
    
    Returns:
        List of (text, slide_metadata) tuples
    """
    # Check file extension
    if file_path.suffix.lower() == ".ppt":
        raise UnsupportedFileFormatError(
            f"Old .ppt format (binary) is not supported. "
            f"Please convert to .pptx format (Office Open XML). "
            f"You can open the file in PowerPoint and 'Save As' .pptx format. "
            f"File: {file_path.name}"
        )
    
    try:
        slides = []
        prs = Presentation(file_path)
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_texts = []
            slide_title = None
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    slide_texts.append(text)
                    # Try to identify title (usually first shape or placeholder)
                    if slide_title is None and shape.shape_type == 1:  # Placeholder type
                        slide_title = text[:100]  # First 100 chars
            
            if slide_texts:
                slide_text = "\n".join(slide_texts)
                slide_metadata = {
                    "slide_index": slide_idx,
                    "slide_title": slide_title or "",
                }
                slides.append((slide_text, slide_metadata))
        
        return slides
    except Exception as e:
        error_msg = str(e)
        if "not a zip file" in error_msg.lower() or "bad zipfile" in error_msg.lower():
            raise DocumentParseError(
                f"File '{file_path}' appears to be in old .ppt format (binary), not .pptx. "
                f"python-pptx only supports .pptx format (Office Open XML). "
                f"Please convert the file to .pptx format (open in PowerPoint and 'Save As' .pptx). "
                f"Original error: {e}"
            ) from e
        raise DocumentParseError(f"Failed to parse PPT: {e}") from e


def parse_document(file_path: Path) -> Document:
    """Parse PDF or PPT file into Document with unit-level structure."""
    file_path = Path(file_path)
    if not file_path.exists():
        raise DocumentParseError(f"File not found: {file_path}")

    suffix = file_path.suffix.lower()
    doc_id = file_path.stem

    base_metadata = {
        "file_path": str(file_path),
        "file_name": file_path.name,
        "file_type": suffix,
    }

    if suffix == ".pdf":
        pages = parse_pdf(file_path)
        # Store pages as list of (text, page_num) for processing
        text = "\n\n---PAGE_SEPARATOR---\n\n".join([f"---PAGE_{page_num}---\n{text}" for text, page_num in pages])
        base_metadata["num_pages"] = len(pages)
        base_metadata["pages"] = [(text, page_num) for text, page_num in pages]
    elif suffix == ".pptx":
        slides = parse_ppt(file_path)
        # Store slides as list of (text, slide_metadata) for processing
        text = "\n\n---SLIDE_SEPARATOR---\n\n".join([f"---SLIDE_{idx}---\n{text}" for idx, (text, _) in enumerate(slides)])
        base_metadata["num_slides"] = len(slides)
        base_metadata["slides"] = slides
    else:
        raise UnsupportedFileFormatError(
            f"Unsupported file format: {suffix}. "
            f"Only .pdf and .pptx formats are supported. "
            f"File: {file_path.name}"
        )

    return Document(doc_id=doc_id, text=text, metadata=base_metadata)


def chunk_document(
    document: Document,
    chunk_size_tokens: int = 1000,
    overlap_sentences: int = 2,
) -> List[Chunk]:
    """Chunk a document using hierarchical recursive splitting with proper metadata."""
    chunks = []
    global_chunk_index = 0
    
    file_type = document.metadata.get("file_type", "")
    
    if file_type == ".pdf":
        pages = document.metadata.get("pages", [])
        for text, page_num in pages:
            # Recursively split this page
            page_chunks = recursive_split(text, chunk_size_tokens)
            
            # Apply sentence overlap
            page_chunks = apply_sentence_overlap(page_chunks, overlap_sentences)
            
            for local_idx, chunk_text in enumerate(page_chunks):
                # Clean any separator markers
                cleaned_text = chunk_text.replace("---PAGE_SEPARATOR---", "").replace(f"---PAGE_{page_num}---", "").strip()
                
                chunk_metadata = {
                    "file_type": file_type,
                    "file_path": document.metadata.get("file_path"),
                    "file_name": document.metadata.get("file_name"),
                    "page_num": page_num,
                    "source_unit": "page",
                    "chunk_local_index": local_idx,
                }
                
                chunk = Chunk(
                    chunk_id=f"{document.doc_id}_page{page_num}_chunk{local_idx}",
                    doc_id=document.doc_id,
                    text=cleaned_text,
                    chunk_index=global_chunk_index,
                    metadata=chunk_metadata,
                )
                chunks.append(chunk)
                global_chunk_index += 1
    
    elif file_type in [".pptx"]:
        slides = document.metadata.get("slides", [])
        for slide_idx, (text, slide_meta) in enumerate(slides):
            # Recursively split this slide
            slide_chunks = recursive_split(text, chunk_size_tokens)
            
            # Apply sentence overlap
            slide_chunks = apply_sentence_overlap(slide_chunks, overlap_sentences)
            
            for local_idx, chunk_text in enumerate(slide_chunks):
                # Clean any separator markers
                cleaned_text = chunk_text.replace("---SLIDE_SEPARATOR---", "").replace(f"---SLIDE_{slide_idx}---", "").strip()
                
                chunk_metadata = {
                    "file_type": file_type,
                    "file_path": document.metadata.get("file_path"),
                    "file_name": document.metadata.get("file_name"),
                    "slide_index": slide_meta.get("slide_index", slide_idx),
                    "slide_title": slide_meta.get("slide_title", ""),
                    "source_unit": "slide",
                    "chunk_local_index": local_idx,
                }
                
                chunk = Chunk(
                    chunk_id=f"{document.doc_id}_slide{slide_idx}_chunk{local_idx}",
                    doc_id=document.doc_id,
                    text=cleaned_text,
                    chunk_index=global_chunk_index,
                    metadata=chunk_metadata,
                )
                chunks.append(chunk)
                global_chunk_index += 1
    
    logger.info(f"Created {len(chunks)} chunks from {document.doc_id} (hierarchical recursive with {overlap_sentences} sentence overlap)")
    return chunks
